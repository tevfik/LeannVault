#!/usr/bin/env python3
"""
LEANN Indexing Script for Office 365 corpus subset.
Extracts text from PDF, PPTX, DOCX, and JSON (email) files,
then indexes them using LEANN with HNSW backend.
"""

import os
import sys
import json
from pathlib import Path
from typing import Optional

PROJECT_ROOT = Path(__file__).parent.parent
DATA_DIR = PROJECT_ROOT / "data" / "subset"
INDEX_PATH = PROJECT_ROOT / "results" / "index.leann"

try:
    from leann import LeannBuilder
except ImportError:
    print("ERROR: LEANN not installed. Run: pip install leann")
    sys.exit(1)

try:
    from pptx import Presentation
except ImportError:
    print("WARNING: python-pptx not installed. PPTX files will be skipped.")
    Presentation = None

try:
    from docx import Document
except ImportError:
    print("WARNING: python-docx not installed. DOCX files will be skipped.")
    Document = None

try:
    import pdfplumber
except ImportError:
    print("WARNING: pdfplumber not installed. PDF files will be skipped.")
    pdfplumber = None


def extract_text_from_pptx(file_path: Path) -> Optional[str]:
    if Presentation is None:
        return None
    try:
        prs = Presentation(str(file_path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n".join(texts) if texts else None
    except Exception as e:
        print(f"  Error reading PPTX {file_path.name}: {e}")
        return None


def extract_text_from_docx(file_path: Path) -> Optional[str]:
    if Document is None:
        return None
    try:
        doc = Document(str(file_path))
        texts = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
        return "\n".join(texts) if texts else None
    except Exception as e:
        print(f"  Error reading DOCX {file_path.name}: {e}")
        return None


def extract_text_from_pdf(file_path: Path) -> Optional[str]:
    if pdfplumber is None:
        return None
    try:
        texts = []
        with pdfplumber.open(str(file_path)) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text and text.strip():
                    texts.append(text.strip())
        return "\n".join(texts) if texts else None
    except Exception as e:
        print(f"  Error reading PDF {file_path.name}: {e}")
        return None


def extract_text_from_json_email(file_path: Path) -> Optional[str]:
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        parts = []
        if "subject" in data:
            parts.append(f"Subject: {data['subject']}")
        if "body" in data:
            parts.append(data["body"])
        elif "bodyPreview" in data:
            parts.append(data["bodyPreview"])
        if "from" in data:
            parts.append(f"From: {data['from']}")
        return "\n".join(parts) if parts else None
    except Exception as e:
        print(f"  Error reading JSON {file_path.name}: {e}")
        return None


def get_extractor(file_path: Path):
    ext = file_path.suffix.lower()
    extractors = {
        ".pptx": extract_text_from_pptx,
        ".docx": extract_text_from_docx,
        ".pdf": extract_text_from_pdf,
        ".json": extract_text_from_json_email,
    }
    return extractors.get(ext)


def scan_and_extract(data_dir: Path):
    documents = []
    file_count = 0
    success_count = 0

    for file_path in data_dir.rglob("*"):
        if not file_path.is_file():
            continue

        extractor = get_extractor(file_path)
        if extractor is None:
            continue

        file_count += 1
        print(f"  Processing ({file_count}): {file_path.name}")

        text = extractor(file_path)
        if text and len(text.strip()) > 50:
            documents.append(
                {
                    "text": text.strip(),
                    "metadata": {
                        "source": str(file_path.relative_to(data_dir)),
                        "type": file_path.suffix.lower(),
                        "size_bytes": file_path.stat().st_size,
                    },
                }
            )
            success_count += 1

    return documents, file_count, success_count


def main():
    print("=" * 60)
    print("LEANN Indexing Script")
    print("=" * 60)

    if not DATA_DIR.exists():
        print(f"ERROR: Data directory not found: {DATA_DIR}")
        sys.exit(1)

    print(f"\nData directory: {DATA_DIR}")
    print(f"Index output: {INDEX_PATH}")
    print()

    print("Scanning and extracting text from documents...")
    documents, total_files, success_count = scan_and_extract(DATA_DIR)

    print(f"\nExtraction summary:")
    print(f"  Files scanned: {total_files}")
    print(f"  Documents extracted: {success_count}")

    if not documents:
        print("ERROR: No documents extracted. Check file types and dependencies.")
        sys.exit(1)

    print(f"\nBuilding LEANN index with HNSW backend...")
    print(f"  Documents to index: {len(documents)}")

    builder = LeannBuilder(backend_name="hnsw")

    for i, doc in enumerate(documents):
        builder.add_text(doc["text"], metadata=doc["metadata"])
        if (i + 1) % 10 == 0:
            print(f"  Indexed {i + 1}/{len(documents)} documents...")

    print(f"  Finalizing index...")
    INDEX_PATH.parent.mkdir(parents=True, exist_ok=True)
    builder.build_index(str(INDEX_PATH))

    index_dir = INDEX_PATH.parent
    index_files = list(index_dir.glob("index.*")) + list(index_dir.glob("*.leann.*"))
    index_size = sum(f.stat().st_size for f in index_files if f.is_file())
    index_size_mb = index_size / (1024 * 1024)

    data_size = sum(f.stat().st_size for f in DATA_DIR.rglob("*") if f.is_file())
    data_size_mb = data_size / (1024 * 1024)

    print(f"\n" + "=" * 60)
    print("INDEXING COMPLETE")
    print("=" * 60)
    print(f"  Index directory: {index_dir}")
    print(f"  Index files: {len(index_files)}")
    print(f"  Index size: {index_size_mb:.2f} MB ({index_size:,} bytes)")
    print(f"  Original data size: {data_size_mb:.2f} MB")
    if data_size > 0:
        print(f"  Compression ratio: {(index_size / data_size * 100):.2f}% of original")
        print(f"  Storage savings: {(100 - index_size / data_size * 100):.2f}%")
    print(f"  Documents indexed: {len(documents)}")


if __name__ == "__main__":
    main()
