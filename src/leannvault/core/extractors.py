"""
Text extractors for various file formats.

Uses Microsoft's markitdown as the primary extraction engine for:
PDF, DOCX, PPTX, XLSX, HTML, and other common formats.

Maintains custom logic for JSON email extraction (Office 365 schema).
"""

from pathlib import Path
from typing import Optional


def extract_with_markitdown(file_path: Path) -> Optional[str]:
    """
    Extract text using Microsoft's markitdown library.

    Supports: PDF, DOCX, PPTX, XLSX, HTML, and many other formats.
    Falls back gracefully if markitdown is not installed or fails.

    Args:
        file_path: Path to the file.

    Returns:
        Extracted text or None if extraction failed.
    """
    try:
        from markitdown import MarkItDown
    except ImportError:
        return None

    try:
        md = MarkItDown()
        result = md.convert(str(file_path))
        if result and result.text_content:
            text = result.text_content.strip()
            return text if text else None
        return None
    except Exception:
        return None


def extract_text_from_pdf(file_path: Path) -> Optional[str]:
    """
    Extract text from a PDF file.

    Uses markitdown as primary, with pdfplumber fallback.

    Args:
        file_path: Path to the PDF file.

    Returns:
        Extracted text or None if extraction failed.
    """
    text = extract_with_markitdown(file_path)
    if text:
        return text

    try:
        import pdfplumber
    except ImportError:
        return None

    try:
        texts = []
        with pdfplumber.open(str(file_path)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    texts.append(page_text.strip())
        return "\n".join(texts) if texts else None
    except Exception:
        return None


def extract_text_from_pptx(file_path: Path) -> Optional[str]:
    """
    Extract text from a PowerPoint file.

    Uses markitdown as primary, with python-pptx fallback.

    Args:
        file_path: Path to the PPTX file.

    Returns:
        Extracted text or None if extraction failed.
    """
    text = extract_with_markitdown(file_path)
    if text:
        return text

    try:
        from pptx import Presentation
    except ImportError:
        return None

    try:
        prs = Presentation(str(file_path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n".join(texts) if texts else None
    except Exception:
        return None


def extract_text_from_docx(file_path: Path) -> Optional[str]:
    """
    Extract text from a Word document.

    Uses markitdown as primary, with python-docx fallback.

    Args:
        file_path: Path to the DOCX file.

    Returns:
        Extracted text or None if extraction failed.
    """
    text = extract_with_markitdown(file_path)
    if text:
        return text

    try:
        from docx import Document
    except ImportError:
        return None

    try:
        doc = Document(str(file_path))
        texts = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
        return "\n".join(texts) if texts else None
    except Exception:
        return None


def extract_text_from_xlsx(file_path: Path) -> Optional[str]:
    """
    Extract text from an Excel file.

    Uses markitdown as primary extraction engine.

    Args:
        file_path: Path to the XLSX file.

    Returns:
        Extracted text or None if extraction failed.
    """
    return extract_with_markitdown(file_path)


def extract_text_from_html(file_path: Path) -> Optional[str]:
    """
    Extract text from an HTML file.

    Uses markitdown as primary extraction engine.

    Args:
        file_path: Path to the HTML file.

    Returns:
        Extracted text or None if extraction failed.
    """
    return extract_with_markitdown(file_path)


def extract_text_from_json_email(file_path: Path) -> Optional[str]:
    """
    Extract text from a JSON email export (Office 365 schema).

    Custom extraction for Office 365 email JSON format.
    Handles subject, body, bodyPreview, and from fields.

    Args:
        file_path: Path to the JSON file.

    Returns:
        Extracted text or None if extraction failed.
    """
    import json

    try:
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        parts = []
        if "subject" in data:
            parts.append(f"Subject: {data['subject']}")
        if "body" in data:
            body = data["body"]
            if isinstance(body, dict) and "content" in body:
                parts.append(body["content"])
            else:
                parts.append(str(body))
        elif "bodyPreview" in data:
            parts.append(data["bodyPreview"])
        if "from" in data:
            from_data = data["from"]
            if isinstance(from_data, dict) and "emailAddress" in from_data:
                email_info = from_data["emailAddress"]
                name = email_info.get("name", "")
                address = email_info.get("address", "")
                parts.append(f"From: {name} <{address}>")
            else:
                parts.append(f"From: {from_data}")
        if "toRecipients" in data:
            recipients = data["toRecipients"]
            if isinstance(recipients, list):
                to_list = []
                for r in recipients:
                    if isinstance(r, dict) and "emailAddress" in r:
                        email_info = r["emailAddress"]
                        name = email_info.get("name", "")
                        address = email_info.get("address", "")
                        to_list.append(f"{name} <{address}>" if name else address)
                if to_list:
                    parts.append(f"To: {', '.join(to_list)}")
        return "\n".join(parts) if parts else None
    except Exception:
        return None


def extract_text(file_path: Path) -> Optional[str]:
    """
    Extract text from a file using the appropriate extractor.

    Routes to the correct extractor based on file extension.
    Uses markitdown for most formats, custom logic for JSON emails.

    Args:
        file_path: Path to the file.

    Returns:
        Extracted text or None if extraction failed.
    """
    ext = file_path.suffix.lower()
    extractors = {
        ".pdf": extract_text_from_pdf,
        ".pptx": extract_text_from_pptx,
        ".ppt": extract_text_from_pptx,
        ".docx": extract_text_from_docx,
        ".doc": extract_text_from_docx,
        ".xlsx": extract_text_from_xlsx,
        ".xls": extract_text_from_xlsx,
        ".html": extract_text_from_html,
        ".htm": extract_text_from_html,
        ".json": extract_text_from_json_email,
    }
    extractor = extractors.get(ext)
    if extractor:
        return extractor(file_path)
    return None


SUPPORTED_EXTENSIONS = (
    set(extract_text.__doc__.split("Routes")[0].split("Uses")[0].strip().split())
    if False
    else {".pdf", ".pptx", ".ppt", ".docx", ".doc", ".xlsx", ".xls", ".html", ".htm", ".json"}
)
